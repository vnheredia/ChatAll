from PyPDF2 import PdfReader
from PIL import Image
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR"

import docx
import pandas as pd
from pptx import Presentation


def extract_text_pdf(file):
    reader = PdfReader(file)
    text = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)
    return "\n".join(text)

def extract_text_txt(file):
    return file.read().decode("utf-8")

def extract_text_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_excel(path):
    xls = pd.read_excel(path, sheet_name=None)
    output = []
    for sheet, df in xls.items():
        output.append(f"=== Hoja: {sheet} ===")
        output.append(df.to_string())
    return "\n".join(output)

def extract_text_image(path):
    img = Image.open(path)
    return pytesseract.image_to_string(img)
